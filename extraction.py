from pptx import Presentation
import requests
import os
def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text_content = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content += shape.text + "\n"
    return text_content
# Mistral API function
MISTRAL_API_URL = "https://api.mistral.ai/v1/chat/completions"
API_KEY = os.getenv("MISTRAL_API_KEY")
def query_mistral_api(prompt):
    headers = {
        "Authorization": f"Bearer {API_KEY}",
        "Content-Type": "application/json",
    }
    payload = {
        "prompt": prompt,
        "max_tokens": 500,
    }
    response = requests.post(MISTRAL_API_URL, json=payload, headers=headers)
    return response.json()
# Full integration: Extract text and query Mistral API
def process_ppt_with_mistral(ppt_file):
    # Extract text from PowerPoint
    text = extract_text_from_pptx(ppt_file)
    # Send extracted text to Mistral API
    prompt = "Summarize the key points from the following PowerPoint content: \n" + text
    result = query_mistral_api(prompt)
    # Return the result
    return result
# Example usage
#ppt_file = "your-presentation.pptx"
#mistral_result = process_ppt_with_mistral(ppt_file)
#print(mistral_result)
